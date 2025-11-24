#!/usr/bin/python3.11
"""
Script per generare PDF da template PowerPoint sostituendo i placeholder.
Uso: python3 pptx-generator.py input.json output.pdf
"""

import sys
import json
import os
import subprocess
import tempfile
from pptx import Presentation
from pathlib import Path

def replace_placeholders_in_shape(shape, replacements):
    """Sostituisce i placeholder in un shape (testo, tabella, gruppo, etc.) preservando la formattazione"""
    from pptx.shapes.group import GroupShape
    
    # Se è un gruppo, processa ricorsivamente tutti gli shape al suo interno
    if isinstance(shape, GroupShape):
        for sub_shape in shape.shapes:
            replace_placeholders_in_shape(sub_shape, replacements)
        return
    
    # Processa text frame
    if shape.has_text_frame:
        replace_in_text_frame(shape.text_frame, replacements)
    
    # Processa tabelle
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                replace_in_text_frame(cell.text_frame, replacements)

def replace_in_text_frame(text_frame, replacements):
    """Sostituisce i placeholder in un text_frame preservando la formattazione"""
    for paragraph in text_frame.paragraphs:
        # Prima raccogli tutto il testo del paragrafo
        full_text = ''.join(run.text for run in paragraph.runs)
        
        # Controlla se ci sono placeholder da sostituire
        needs_replacement = any(placeholder in full_text for placeholder in replacements.keys())
        
        if needs_replacement:
            # Sostituisci i placeholder
            new_text = full_text
            for placeholder, value in replacements.items():
                new_text = new_text.replace(placeholder, value)
            
            # Se il testo è cambiato, aggiorna il primo run e rimuovi gli altri
            if new_text != full_text:
                # Mantieni solo il primo run con tutto il nuovo testo
                if paragraph.runs:
                    paragraph.runs[0].text = new_text
                    # Rimuovi gli altri run
                    for i in range(len(paragraph.runs) - 1, 0, -1):
                        paragraph._element.remove(paragraph.runs[i]._element)

def generate_pdf_from_template(data, template_path, output_pdf_path):
    """
    Genera PDF da template PowerPoint sostituendo i placeholder.
    
    Args:
        data: Dict con i dati della valutazione
        template_path: Path al template .pptx
        output_pdf_path: Path dove salvare il PDF finale
    """
    
    # Prepara i replacements (usa chiavi UPPERCASE come in routers.ts)
    replacements = {
        '{{COMUNE}}': data.get('COMUNE', ''),
        '{{LOCALITA}}': data.get('LOCALITA', ''),
        '{{VALORE_TOTALE}}': data.get('VALORE_TOTALE', '0'),
        '{{VALORE_MINIMO}}': data.get('VALORE_MINIMO', '0'),
        '{{VALORE_MASSIMO}}': data.get('VALORE_MASSIMO', '0'),
        '{{PREZZO_MQ}}': data.get('PREZZO_MQ', '0'),
        '{{COMPETITIVITA}}': data.get('COMPETITIVITA', ''),
        '{{PREZZO_CONSIGLIATO}}': data.get('PREZZO_CONSIGLIATO', '0'),
        '{{TIPOLOGIA}}': data.get('TIPOLOGIA', ''),
        '{{PIANO}}': data.get('PIANO', ''),
        '{{STATO}}': data.get('STATO', ''),
        '{{VISTA_MARE}}': data.get('VISTA_MARE', ''),
        '{{DISTANZA_MARE}}': data.get('DISTANZA_MARE', ''),
        '{{SUPERFICIE}}': data.get('SUPERFICIE', '0 mq'),
        '{{VALORE_BASE}}': data.get('VALORE_BASE', '0'),
        '{{VALORE_PERTINENZE}}': data.get('VALORE_PERTINENZE', '0'),
        '{{VALORE_VALORIZZAZIONI}}': data.get('VALORE_VALORIZZAZIONI', '0'),
        
        # Punti di forza dinamici
        '{{ICONA_1}}': data.get('ICONA_1', '🌊'),
        '{{PUNTO_FORZA_1_TITOLO}}': data.get('PUNTO_FORZA_1_TITOLO', ''),
        '{{PUNTO_FORZA_1_TESTO}}': data.get('PUNTO_FORZA_1_TESTO', ''),
        
        '{{ICONA_2}}': data.get('ICONA_2', '🏖️'),
        '{{PUNTO_FORZA_2_TITOLO}}': data.get('PUNTO_FORZA_2_TITOLO', ''),
        '{{PUNTO_FORZA_2_TESTO}}': data.get('PUNTO_FORZA_2_TESTO', ''),
        
        '{{ICONA_3}}': data.get('ICONA_3', '⭐'),
        '{{PUNTO_FORZA_3_TITOLO}}': data.get('PUNTO_FORZA_3_TITOLO', ''),
        '{{PUNTO_FORZA_3_TESTO}}': data.get('PUNTO_FORZA_3_TESTO', ''),
        
        '{{ICONA_4}}': data.get('ICONA_4', '🌳'),
        '{{PUNTO_FORZA_4_TITOLO}}': data.get('PUNTO_FORZA_4_TITOLO', ''),
        '{{PUNTO_FORZA_4_TESTO}}': data.get('PUNTO_FORZA_4_TESTO', ''),
    }
    
    # Carica il template PowerPoint
    prs = Presentation(template_path)
    
    # Sostituisci i placeholder in tutte le slide
    for slide in prs.slides:
        for shape in slide.shapes:
            replace_placeholders_in_shape(shape, replacements)
    
    # Salva il PowerPoint modificato in un file temporaneo
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp_pptx:
        tmp_pptx_path = tmp_pptx.name
        prs.save(tmp_pptx_path)
    
    try:
        # Converti PPTX in PDF usando LibreOffice
        convert_pptx_to_pdf(tmp_pptx_path, output_pdf_path)
    finally:
        # Rimuovi il file temporaneo
        os.unlink(tmp_pptx_path)

def convert_pptx_to_pdf(pptx_path, pdf_path):
    """Converte PPTX in PDF usando LibreOffice headless"""
    output_dir = os.path.dirname(pdf_path)
    
    # Esegui LibreOffice in modalità headless per conversione
    cmd = [
        'libreoffice',
        '--headless',
        '--convert-to', 'pdf',
        '--outdir', output_dir,
        pptx_path
    ]
    
    result = subprocess.run(cmd, capture_output=True, text=True)
    
    if result.returncode != 0:
        raise Exception(f"Errore conversione PDF: {result.stderr}")
    
    # LibreOffice salva con il nome del file input, rinomina se necessario
    generated_pdf = os.path.join(output_dir, os.path.basename(pptx_path).replace('.pptx', '.pdf'))
    if generated_pdf != pdf_path and os.path.exists(generated_pdf):
        os.rename(generated_pdf, pdf_path)

def format_currency(value):
    """Formatta un numero come valuta italiana"""
    if isinstance(value, str):
        return value
    return f"{int(value):,}".replace(',', '.')

def main():
    if len(sys.argv) != 3:
        print("Uso: python3 pptx-generator.py input.json output.pdf")
        sys.exit(1)
    
    input_json = sys.argv[1]
    output_pdf = sys.argv[2]
    
    # Leggi i dati dal JSON
    with open(input_json, 'r', encoding='utf-8') as f:
        data = json.load(f)
    
    # Path al template
    script_dir = os.path.dirname(os.path.abspath(__file__))
    template_path = os.path.join(script_dir, 'templates', 'template-stima.pptx')
    
    # Genera il PDF
    generate_pdf_from_template(data, template_path, output_pdf)
    
    print(f"✅ PDF generato: {output_pdf}")

if __name__ == '__main__':
    main()
